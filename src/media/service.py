"""Media upload service — stores files to S3 or local disk, extracts text, and indexes chunks in pgvector."""
import uuid
import os
import asyncio
from pathlib import Path
from typing import Optional, Dict, Any
from src.common.logger import get_logger
from src.config import settings

logger = get_logger(__name__)


async def upload_file(
    file_data: bytes,
    filename: str,
    content_type: str,
    uploader_id: str,
    group_id: Optional[str] = None,
    description: str = "",
) -> Dict[str, Any]:
    media_id = str(uuid.uuid4())
    ext = Path(filename).suffix
    stored_name = f"{media_id}{ext}"

    url = await _try_s3_upload(file_data, stored_name, content_type)
    if not url:
        url = await _local_upload(file_data, stored_name)

    media_type = _infer_media_type(content_type)
    record = {
        "media_id": media_id,
        "original_name": filename,
        "stored_name": stored_name,
        "content_type": content_type,
        "media_type": media_type,
        "url": url,
        "size_bytes": len(file_data),
        "uploader_id": uploader_id,
        "group_id": group_id,
        "description": description,
    }

    from src.common.database import get_mongo_db
    db = get_mongo_db()
    await db.media.insert_one({**record})

    try:
        from src.ai.vector_store import get_vector_store
        vs = get_vector_store()
        if vs:
            embed_text = _extract_text(file_data, content_type, filename) or description or filename
            await vs.add_document_chunks(media_id, embed_text, {
                "filename": filename,
                "media_type": media_type,
                "uploader_id": uploader_id,
                "group_id": group_id or "",
                "receiver_id": "",
            })
    except Exception as e:
        logger.warning("media_index_failed", error=str(e))

    logger.info("media_uploaded", media_id=media_id, type=media_type)
    return record


async def _try_s3_upload(data: bytes, key: str, content_type: str) -> Optional[str]:
    if not settings.AWS_ACCESS_KEY_ID:
        return None
    try:
        import boto3
        s3 = boto3.client(
            "s3",
            aws_access_key_id=settings.AWS_ACCESS_KEY_ID,
            aws_secret_access_key=settings.AWS_SECRET_ACCESS_KEY,
            region_name=settings.AWS_REGION,
        )
        loop = asyncio.get_event_loop()
        await loop.run_in_executor(
            None,
            lambda: s3.put_object(
                Bucket=settings.AWS_S3_BUCKET,
                Key=key,
                Body=data,
                ContentType=content_type,
                # ACL omitted — use bucket policy for public read instead of
                # per-object ACL (Block Public Access rejects ACL on new buckets)
            ),
        )
        url = f"https://{settings.AWS_S3_BUCKET}.s3.{settings.AWS_REGION}.amazonaws.com/{key}"
        logger.info("s3_upload_ok", key=key, url=url)
        return url
    except Exception as e:
        logger.warning("s3_upload_failed", error=str(e), bucket=settings.AWS_S3_BUCKET, key=key)
        return None


async def _local_upload(data: bytes, filename: str) -> str:
    upload_dir = Path(settings.LOCAL_UPLOADS_PATH)
    upload_dir.mkdir(parents=True, exist_ok=True)
    path = upload_dir / filename
    loop = asyncio.get_event_loop()
    await loop.run_in_executor(None, path.write_bytes, data)
    return f"/uploads/{filename}"


def _infer_media_type(content_type: str) -> str:
    if content_type.startswith("image/"):
        return "image"
    if content_type.startswith("video/"):
        return "video"
    if content_type.startswith("audio/"):
        return "voice"
    return "document"


def _extract_text(data: bytes, content_type: str, filename: str) -> str:
    try:
        if content_type == "application/pdf":
            import fitz
            doc = fitz.open(stream=data, filetype="pdf")
            return "\n".join(page.get_text() for page in doc)[:8000]
        if "word" in content_type or filename.lower().endswith((".docx", ".doc")):
            from docx import Document
            from io import BytesIO
            doc = Document(BytesIO(data))
            return "\n".join(p.text for p in doc.paragraphs if p.text.strip())[:8000]
        if "presentation" in content_type or filename.lower().endswith((".pptx", ".ppt")):
            from pptx import Presentation
            from io import BytesIO
            prs = Presentation(BytesIO(data))
            texts = [
                shape.text
                for slide in prs.slides
                for shape in slide.shapes
                if hasattr(shape, "text") and shape.text.strip()
            ]
            return "\n".join(texts)[:8000]
        if content_type.startswith("text/"):
            return data.decode("utf-8", errors="replace")[:8000]
    except Exception as e:
        logger.warning("text_extraction_failed", content_type=content_type, error=str(e))
    return ""
